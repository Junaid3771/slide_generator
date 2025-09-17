from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO
from PIL import Image as PILImage
from typing import Union
from typing import Union, List
from pptx.oxml import parse_xml
from pptx.enum.text import MSO_AUTO_SIZE
def create_title_slide(
    prs: Presentation,
    title: str,
    subtitle: str = "",
    content: str = "",
    image: Union[PILImage.Image, bytes, None] = None
) -> Presentation:
    """
    Add a title-style slide to `prs`:
      - Left: title / subtitle / content (stacked, no overlap)
      - Right bottom corner: image in a styled container
    """

    # --- Create blank slide ---
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # --- Gradient background ---
    background = slide.background
    fill = background.fill
    try:
        fill.gradient()
        grad_stops = fill.gradient_stops
        grad_stops[0].color.rgb = RGBColor(255, 230, 230)
        grad_stops[1].color.rgb = RGBColor(255, 255, 255)
    except Exception:
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 250, 250)

    # --- Geometry setup ---
    margin = Inches(0.7)
    padding_between = Inches(0.25)

    slide_w, slide_h = prs.slide_width, prs.slide_height

    # Left column (text)
    text_col_left = margin
    text_col_width = slide_w - margin - Inches(4.5) - padding_between  # leave space for image

    # --- Title ---
    title_box = slide.shapes.add_textbox(text_col_left, Inches(0.9), text_col_width, Inches(2))
    title_tf = title_box.text_frame
    title_tf.clear()
    p = title_tf.paragraphs[0]
    p.text = title
    r = p.runs[0]
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(10, 10, 10)
    r.font.name = "Segoe UI"

    title_tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    title_bottom = title_box.top + title_box.height + Pt(6)

    # --- Subtitle ---
    if subtitle:
        sub_box = slide.shapes.add_textbox(text_col_left, title_bottom, text_col_width, Inches(1))
        sub_tf = sub_box.text_frame
        sub_tf.clear()
        p = sub_tf.paragraphs[0]
        p.text = subtitle
        r = p.runs[0]
        r.font.size = Pt(18)
        r.font.italic = True
        r.font.color.rgb = RGBColor(95, 95, 95)
        r.font.name = "Segoe UI"
        subtitle_bottom = sub_box.top + sub_box.height + Pt(6)
    else:
        subtitle_bottom = title_bottom

    # --- Content ---
    if content:
        content_h = slide_h - subtitle_bottom - margin
        content_box = slide.shapes.add_textbox(text_col_left, subtitle_bottom, text_col_width, content_h)
        content_tf = content_box.text_frame
        content_tf.clear()
        p = content_tf.paragraphs[0]
        p.text = content
        r = p.runs[0]
        r.font.size = Pt(14)
        r.font.color.rgb = RGBColor(70, 70, 70)
        r.font.name = "Segoe UI"
        content_tf.word_wrap = True
        for para in content_tf.paragraphs:
            para.space_after = Pt(6)

    # --- Image area (bottom-right corner) ---
    image_area_width = Inches(4.0)
    image_area_height = Inches(3.0)
    image_left = slide_w - image_area_width - margin
    image_top = slide_h - image_area_height - margin

    # Decorative container
    container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, image_left, image_top, image_area_width, image_area_height
    )
    container.fill.gradient()
    grad_stops = container.fill.gradient_stops
    grad_stops[0].color.rgb = RGBColor(245, 245, 255)
    grad_stops[1].color.rgb = RGBColor(225, 225, 240)
    container.line.color.rgb = RGBColor(180, 180, 200)
    try:
        container.shadow.inherit = False
        container.shadow.blur_radius = Pt(8)
        container.shadow.distance = Pt(4)
    except Exception:
        pass

    # --- Add image inside container ---
    if image:
        try:
            if isinstance(image, bytes):
                image_obj = PILImage.open(BytesIO(image))
            elif isinstance(image, PILImage.Image):
                image_obj = image
            else:
                image_obj = PILImage.open(image)

            orig_w, orig_h = image_obj.size
            target_w = int(image_area_width * 0.9)
            scaled_h = int(target_w * (orig_h / orig_w))

            if scaled_h > (image_area_height * 0.9):
                target_h = int(image_area_height * 0.9)
                scaled_w = int(target_h * (orig_w / orig_h))
                target_w, scaled_h = scaled_w, target_h

            img_left = image_left + int((image_area_width - target_w) / 2)
            img_top = image_top + int((image_area_height - scaled_h) / 2)

            img_stream = BytesIO()
            if image_obj.mode in ("RGBA", "P"):
                image_obj = image_obj.convert("RGB")
            image_obj.save(img_stream, format="PNG")
            img_stream.seek(0)

            slide.shapes.add_picture(img_stream, img_left, img_top, width=target_w, height=scaled_h)

        except Exception:
            _add_image_placeholder(slide, image_left, image_top, image_area_width, image_area_height, "Image Error")
    else:
        _add_image_placeholder(slide, image_left, image_top, image_area_width, image_area_height)

    return prs


def _add_image_placeholder(slide, left, top, width, height, text: str):
    """Helper: add a placeholder rectangle with text."""
    placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
    placeholder.line.color.rgb = RGBColor(200, 200, 200)
    tf = placeholder.text_frame
    tf.text = text
    tf.paragraphs[0].runs[0].font.size = Pt(14)
    tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(120, 120, 120)





def create_bullet_slide(
    prs: Presentation,
    title: str,
    bullet_points: List[str],
    image: Union[PILImage.Image, bytes, None] = None
) -> Presentation:
    """
    Add a bullet-point slide to `prs`:
      - Left: Title and bullet points
      - Right: Image inside a styled container (auto-scaled)
    """

    # --- Create blank slide ---
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # --- Gradient background ---
    background = slide.background
    fill = background.fill
    try:
        fill.gradient()
        grad_stops = fill.gradient_stops
        grad_stops[0].color.rgb = RGBColor(255, 230, 230)
        grad_stops[1].color.rgb = RGBColor(255, 255, 255)
    except Exception:
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 250, 250)

    # --- Geometry ---
    margin = Inches(0.7)
    padding_between = Inches(0.3)
    slide_w, slide_h = prs.slide_width, prs.slide_height

    # Right column (image)
    max_image_col = int(slide_w * 0.42)
    min_image_col = Inches(3.2)
    image_col_width = max(min_image_col, max_image_col)

    # Left column (text)
    text_col_left = margin
    text_col_width = slide_w - margin - image_col_width - padding_between - margin

    # --- Title ---
    title_box = slide.shapes.add_textbox(text_col_left, Inches(0.7), text_col_width, Inches(1.5))  ### CHANGED
    title_tf = title_box.text_frame
    title_tf.clear()
    p = title_tf.paragraphs[0]
    p.text = title
    r = p.runs[0]
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = RGBColor(10, 10, 10)
    r.font.name = "Segoe UI"

    # --- Measure title height ---   ### CHANGED
    title_height = title_box.height
    bullets_top = title_box.top + title_height + Pt(6)   ### CHANGED
    bullets_h = (slide_h / 2) - bullets_top              ### CHANGED (limit to half slide)

    # --- Bullet points ---
    bullets_box = slide.shapes.add_textbox(text_col_left, bullets_top, text_col_width, bullets_h)
    bullets_tf = bullets_box.text_frame
    bullets_tf.clear()
    bullets_tf.word_wrap = True

    for idx, point in enumerate(bullet_points):
        p = bullets_tf.add_paragraph() if idx > 0 else bullets_tf.paragraphs[0]
        p.text = point
        p.level = 0
        # ✅ enforce bullets
        p._element.get_or_add_pPr().insert(
            0,
            parse_xml("<a:buChar xmlns:a='http://schemas.openxmlformats.org/drawingml/2006/main' char='•'/>")
        )
        r = p.runs[0]
        r.font.size = Pt(16)
        r.font.color.rgb = RGBColor(60, 60, 60)
        r.font.name = "Segoe UI"
        p.space_after = Pt(6)
        
        p._element.get_or_add_pPr().set("marL", "342900")   # ~0.3 inch left margin
        p._element.get_or_add_pPr().set("indent", "-171450")  # hanging indent

    # --- Image area (styled container) ---
    image_left = slide_w - image_col_width - margin
    image_area_top = title_box.top + title_height + Inches(0.3)   ### CHANGED (starts below title)
    image_area_width = image_col_width
    image_area_height = slide_h - image_area_top - margin         ### CHANGED (fits rest of slide)

    container = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, image_left, image_area_top, image_area_width, image_area_height
    )
    container.fill.solid()
    container.fill.fore_color.rgb = RGBColor(250, 250, 250)
    container.line.color.rgb = RGBColor(200, 200, 200)

    try:
        container.shadow.inherit = False
        container.shadow.blur_radius = Pt(6)
        container.shadow.distance = Pt(4)
    except Exception:
        pass

    if image:
        try:
            if isinstance(image, bytes):
                image_obj = PILImage.open(BytesIO(image))
            elif isinstance(image, PILImage.Image):
                image_obj = image
            else:
                image_obj = PILImage.open(image)

            orig_w, orig_h = image_obj.size
            target_w = int(image_area_width * 0.85)
            scaled_h = int(target_w * (orig_h / orig_w))

            if scaled_h > (image_area_height * 0.85):
                target_h = int(image_area_height * 0.85)
                scaled_w = int(target_h * (orig_w / orig_h))
                target_w, scaled_h = scaled_w, target_h

            image_left_adj = int(image_left + (image_area_width - target_w) / 2)
            image_top_adj = int(image_area_top + (image_area_height - scaled_h) / 2)

            img_stream = BytesIO()
            if image_obj.mode in ("RGBA", "P"):
                image_obj = image_obj.convert("RGB")
            image_obj.save(img_stream, format="PNG")
            img_stream.seek(0)

            slide.shapes.add_picture(img_stream, image_left_adj, image_top_adj, width=target_w, height=scaled_h)

        except Exception:
            _add_image_placeholder(slide, image_left, image_area_top, image_area_width, image_area_height, "Image error")
    else:
        _add_image_placeholder(slide, image_left, image_area_top, image_area_width, image_area_height, "Image Placeholder")

    return prs

def create_two_column_slide(
    prs: Presentation,
    title: str,
    left_column: List[str],
    right_column: List[str],
    image: Union[PILImage.Image, bytes, None] = None
) -> Presentation:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    # --- Gradient background ---
    background = slide.background
    fill = background.fill
    try:
        fill.gradient()
        grad_stops = fill.gradient_stops
        grad_stops[0].color.rgb = RGBColor(255, 230, 230)
        grad_stops[1].color.rgb = RGBColor(255, 255, 255)
    except Exception:
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 250, 250)


    # --- Setup dimensions ---
    margin = Inches(0.7)
    slide_w, slide_h = prs.slide_width, prs.slide_height

    # Title box
    title_box = slide.shapes.add_textbox(margin, Inches(0.5), slide_w - 2 * margin, Inches(1.2))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].runs[0].font.size = Pt(26)
    title_tf.paragraphs[0].runs[0].font.bold = True
    title_tf.paragraphs[0].alignment = 1

    # Content area
    content_top = title_box.top + title_box.height + Inches(0.2)
    content_height = slide_h - content_top - margin  

    # Reserve potential image space (max 3 inches)
    max_image_height = Inches(3)
    min_text_area = Inches(3)  # minimum height we want for text

    # If text area already small, skip image
    if content_height <= min_text_area:
        image = None  

    # Adjust content height depending on whether we keep image
    if image:
        content_height = content_height - max_image_height - Inches(0.3)

    # --- Two-column text boxes ---
    col_width = (slide_w - 2 * margin - Inches(0.4)) / 2
    left_box = slide.shapes.add_textbox(margin, content_top, col_width, content_height)
    right_box = slide.shapes.add_textbox(margin + col_width + Inches(0.4), content_top, col_width, content_height)

    for box, points in [(left_box, left_column), (right_box, right_column)]:
        tf = box.text_frame
        tf.clear()
        for idx, point in enumerate(points):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.text = point
            p.level = 0
            p.runs[0].font.size = Pt(16)

    # --- Image placement (only if enough room left) ---
    if image:
        remaining_space = slide_h - (content_top + content_height) - margin
        if remaining_space >= Inches(2):  # only place if >=2 inches free
            image_area_height = min(remaining_space, max_image_height)
            image_left = margin
            image_top = slide_h - image_area_height - margin
            image_area_width = slide_w - 2 * margin

            # Scale and center image
            try:
                if isinstance(image, bytes):
                    img = PILImage.open(BytesIO(image))
                elif isinstance(image, PILImage.Image):
                    img = image
                else:
                    img = PILImage.open(image)

                orig_w, orig_h = img.size
                scale = min(image_area_width / orig_w, image_area_height / orig_h) * 0.9
                target_w, target_h = int(orig_w * scale), int(orig_h * scale)
                img_left = int(image_left + (image_area_width - target_w) / 2)
                img_top = int(image_top + (image_area_height - target_h) / 2)

                stream = BytesIO()
                img.convert("RGB").save(stream, format="PNG")
                stream.seek(0)
                slide.shapes.add_picture(stream, img_left, img_top, width=target_w, height=target_h)
            except Exception:
                pass  # silently fail if image can't be placed

    return prs
